"""
マルチフォーマット ファイルパーサー (PHASE 3 / Eval-MCP-Phase 1)

対応: PDF / TXT / MD / DOCX / HTML / PPTX / XLSX / CSV / 画像
日本語・英語両対応

【設計方針】
- lxml はオプション。未インストール時は html.parser (Python 標準) を使用
- PDF は pypdf (MIT License) を 50 ページバッチで処理。
  PyMuPDF (AGPL-3.0) からの差し替えで OSS 公開時のライセンス汚染を回避。
- 全パーサーは例外をキャッチしてエラー情報を返す (Publish 全体を止めない)

注: rag.py に既存の取り込み処理がある (._read_file_smart 等)。
本モジュールは独立した補助 API として提供し、Publish パイプラインを
壊さないよう既存処理は維持する。
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Callable, Optional

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".txt",
    ".md",
    ".markdown",
    ".docx",
    ".doc",
    ".html",
    ".htm",
    ".pptx",
    ".ppt",
    ".xlsx",
    ".xls",
    ".csv",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".gif",
    ".bmp",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}


def is_supported(file_path: str) -> bool:
    return Path(file_path).suffix.lower() in SUPPORTED_EXTENSIONS


def is_image(file_path: str) -> bool:
    return Path(file_path).suffix.lower() in IMAGE_EXTENSIONS


def parse_file(
    file_path: str,
    progress_callback: Optional[Callable[[int, int], None]] = None,
    page_batch_size: int = 50,
) -> dict:
    """ファイルを解析してテキストとメタデータを返す。

    Returns:
        {
            "text": str,
            "pages": list[str],  # PDF 用ページ別テキスト
            "page_count": int,
            "file_type": str,
            "metadata": dict,
            "is_image": bool,
            "error": str or None,
        }
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    result = {
        "text": "",
        "pages": [],
        "page_count": 1,
        "file_type": ext,
        "metadata": {"filename": path.name, "path": str(file_path)},
        "is_image": ext in IMAGE_EXTENSIONS,
        "error": None,
    }
    try:
        if ext == ".pdf":
            result.update(_parse_pdf(file_path, progress_callback, page_batch_size))
        elif ext in {".html", ".htm"}:
            result.update(_parse_html(file_path))
        elif ext in {".txt", ".md", ".markdown", ".csv"}:
            result.update(_parse_text(file_path))
        elif ext in {".docx", ".doc"}:
            result.update(_parse_docx(file_path))
        elif ext in {".pptx", ".ppt"}:
            result.update(_parse_pptx(file_path))
        elif ext in {".xlsx", ".xls"}:
            result.update(_parse_xlsx(file_path))
        elif ext in IMAGE_EXTENSIONS:
            result["is_image"] = True  # mlx-vlm 等で別途処理
        else:
            result["error"] = f"Unsupported: {ext}"
    except Exception as e:
        result["error"] = str(e)
    return result


def _parse_pdf(file_path, progress_callback=None, batch_size=50):
    """pypdf (MIT) を使用。PyMuPDF (AGPL-3.0) から差替えで OSS 公開対応。
    1 ページずつ抽出し 50 ページごとに progress_callback を呼ぶ。"""
    import pypdf

    reader = pypdf.PdfReader(file_path)
    page_count = len(reader.pages)
    pages = []
    for batch_start in range(0, page_count, batch_size):
        batch_end = min(batch_start + batch_size, page_count)
        for i in range(batch_start, batch_end):
            try:
                text = reader.pages[i].extract_text() or ""
            except Exception:
                text = ""
            pages.append(text)
        if progress_callback:
            try:
                progress_callback(batch_end, page_count)
            except Exception:
                pass
    return {
        "text": "\n\n".join(p for p in pages if p.strip()),
        "pages": pages,
        "page_count": page_count,
        "metadata": {"page_count": page_count},
    }


def _parse_html(file_path):
    from bs4 import BeautifulSoup

    with open(file_path, "rb") as f:
        raw = _decode_bytes(f.read())
    # lxml が使えれば lxml、なければ html.parser (Python 標準)
    try:
        soup = BeautifulSoup(raw, "lxml")
    except Exception:
        soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(
        [
            "script",
            "style",
            "nav",
            "footer",
            "header",
            "aside",
            "noscript",
            "iframe",
            "form",
            "button",
        ]
    ):
        tag.decompose()
    title = soup.find("title")
    text = soup.get_text(separator="\n", strip=True)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return {
        "text": text,
        "metadata": {"title": title.get_text(strip=True) if title else ""},
    }


def _parse_text(file_path):
    with open(file_path, "rb") as f:
        raw = f.read()
    text = _decode_bytes(raw)
    return {"text": text}


def _decode_bytes(raw: bytes) -> str:
    """日本語ファイルの一般的なエンコーディングを順に試行する。
    UTF-8 BOM, UTF-8, Shift-JIS, CP932, EUC-JP の順で decode 成功した結果を返す。
    最後に utf-8 errors=ignore でフォールバック。"""
    if raw.startswith(b"\xef\xbb\xbf"):
        try:
            return raw.decode("utf-8-sig")
        except UnicodeDecodeError:
            pass
    for enc in ("utf-8", "cp932", "shift_jis", "euc-jp"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _parse_docx(file_path):
    try:
        import docx  # python-docx

        doc = docx.Document(file_path)
        return {"text": "\n".join(p.text for p in doc.paragraphs if p.text.strip())}
    except ImportError:
        return {"text": "", "error": "python-docx not installed"}


def _parse_pptx(file_path):
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        texts = [
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text and shape.text.strip()
        ]
        return {"text": "\n".join(texts), "page_count": len(prs.slides)}
    except ImportError:
        return {"text": "", "error": "python-pptx not installed"}


def _parse_xlsx(file_path):
    try:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                r = "\t".join(str(c) for c in row if c is not None)
                if r.strip():
                    rows.append(r)
        return {"text": "\n".join(rows)}
    except ImportError:
        return {"text": "", "error": "openpyxl not installed"}
